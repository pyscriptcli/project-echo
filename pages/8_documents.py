# pages/8_documents.py
"""
Documents — template-based document generator (native to Project Echo).

Upload or load PPTX/DOCX templates, map {{placeholders}} to Text / Image / Map
fields, auto-fill CTA presets from the advisor contacts database, and download
the generated documents. Includes an interactive map editor that stitches a
high-resolution static map with a draggable pin for Map placeholders.

Native to the app: require_login(), shared sidebar/theme, stone/navy tokens.

Repository templates live in `documents_template/` (committed, protected from
app deletion); templates uploaded through the app are saved to
`stored_templates/` (gitignored).
"""
import os
import sys

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

import io
import re
import json
import math
import time
import streamlit as st
from pptx import Presentation
from PIL import Image, ImageDraw
from datetime import datetime
from docx import Document

# --- MAP SPECIFIC DEPENDENCIES ---
import folium
from folium.plugins import Draw
from streamlit_folium import st_folium
import requests

from utils.auth import require_login
from components.sidebar import setup_page_layout

# ---------------------------------------------------------------------------
# Page shell (native Echo conventions)
# ---------------------------------------------------------------------------
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATES_DIR = os.path.join(PROJECT_ROOT, "documents_template")
STORAGE_DIR = os.path.join(PROJECT_ROOT, "stored_templates")

st.set_page_config(page_title="Documents", layout="wide", initial_sidebar_state="expanded")
require_login()
setup_page_layout()

# --- NATIVE ECHO PAGE CSS (stone/navy tokens, flat & edgy) ---
DOCUMENTS_PAGE_CSS = """
<style>
/* Documents page — custom classes styled with Project Echo tokens */
.editor-card {
    background-color: rgba(249, 250, 251, 0.95);
    border: 2px solid #003366;
    border-radius: 0;
    padding: 1.25rem;
}
.section-header {
    font-family: 'Cormorant Garamond', serif;
    font-style: italic;
    font-weight: 600;
    font-size: 1.5rem;
    color: #003366;
    border-left: 4px solid #003366;
    padding-left: 0.6rem;
    margin-bottom: 0.6rem;
}
.docs-title {
    font-family: 'Cormorant Garamond', serif;
    font-style: italic;
    font-weight: 600;
    font-size: 2.4rem;
    color: #003366;
    line-height: 1.1;
    margin: 0 0 0.2rem 0;
}
.docs-caption {
    font-family: 'Montserrat', sans-serif;
    font-size: 0.85rem;
    color: #69727d;
    margin: 0 0 0.5rem 0;
}
.saved-indicator {
    background-color: #F9FAFB;
    border-left: 3px solid #003366;
    color: #003366;
    font-family: 'Montserrat', sans-serif;
    font-size: 0.8rem;
    font-weight: 600;
    padding: 0.4rem 0.7rem;
    margin-top: 0.4rem;
}
.local-only-note {
    background-color: rgba(249, 250, 251, 0.9);
    border: 1px solid rgba(0,51,102,0.15);
    border-left: 3px solid #003366;
    color: #1b1d1e;
    font-family: 'Montserrat', sans-serif;
    font-size: 0.8rem;
    font-weight: 500;
    padding: 0.45rem 0.7rem;
    margin-bottom: 0.75rem;
}
.placeholder-label,
.field-label {
    font-family: 'Montserrat', sans-serif;
    font-weight: 600;
    font-size: 0.8rem;
    color: #003366;
    margin-bottom: 0.25rem;
    padding-top: 0.35rem;
}
hr {
    margin: 0.75rem 0 !important;
    border-color: rgba(0,51,102,0.15) !important;
}
</style>
"""

st.markdown(DOCUMENTS_PAGE_CSS, unsafe_allow_html=True)

# --- CONTACTS DATABASE FOR CTA PRESETS ---
contacts_database = {
    "Sondi Tuazon": {"phone": "0917 843 6128", "email": "sondi.tuazon@primephilippines.com"},
    "Meliza Zapata": {"phone": "0996 880 5399", "email": "meliza.zapata@primephilippines.com"},
    "Dykstra Pineda": {"phone": "0920 986 2748", "email": "dykstra.pineda@primephilippines.com"},
    "Cedtrix Rena": {"phone": "0977 653 1494", "email": "cedtriz.rena@primephilippines.com"},
    "Carlo Medina": {"phone": "0920 986 2763", "email": "carlo.medina@primephilippines.com"},
    "Dave Policarpio": {"phone": "0908 865 8945", "email": "dave.policarpio@primephilippines.com"},
    "Irish Rima": {"phone": "0918 622 5346", "email": "irish.rima@primephilippines.com"}
}

# --- FILE MANAGEMENT FUNCTIONS ---
def get_storage_dir():
    os.makedirs(STORAGE_DIR, exist_ok=True)
    return STORAGE_DIR

def get_temp_config_path(template_name):
    storage_dir = get_storage_dir()
    safe_name = re.sub(r'[^\w\-_. ]', '_', template_name or "unsaved_template")
    return os.path.join(storage_dir, f"{safe_name}_temp_form_data.json")

def get_github_templates():
    templates = []
    if os.path.exists(TEMPLATES_DIR):
        for file in os.listdir(TEMPLATES_DIR):
            if file.startswith('template_') and (file.endswith('.pptx') or file.endswith('.docx')):
                filepath = os.path.join(TEMPLATES_DIR, file)
                stat = os.stat(filepath)
                display_name = file.replace('template_', '').replace('.pptx', '').replace('.docx', '')
                templates.append({
                    'name': file,
                    'display_name': display_name,
                    'path': filepath,
                    'size': stat.st_size,
                    'modified': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                    'type': 'PPTX' if file.endswith('.pptx') else 'DOCX',
                    'source': 'github'
                })
    return templates

def save_template_to_file(template_bytes, template_name):
    storage_dir = get_storage_dir()
    safe_name = re.sub(r'[^\w\-_. ]', '_', template_name)
    if not safe_name.endswith('.pptx') and not safe_name.endswith('.docx'):
        safe_name += '.docx'
    filepath = os.path.join(storage_dir, safe_name)
    with open(filepath, 'wb') as f:
        f.write(template_bytes)
    get_saved_templates.clear()
    load_template_from_file.clear()
    return filepath

@st.cache_data(show_spinner=False, ttl=60)
def load_template_from_file(template_name):
    templates_filepath = os.path.join(TEMPLATES_DIR, template_name)
    if os.path.exists(templates_filepath):
        with open(templates_filepath, 'rb') as f:
            return f.read()
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, template_name)
    if os.path.exists(filepath):
        with open(filepath, 'rb') as f:
            return f.read()
    return None

@st.cache_data(show_spinner=False, ttl=30)
def get_saved_templates():
    storage_dir = get_storage_dir()
    templates = []
    if os.path.exists(storage_dir):
        for file in os.listdir(storage_dir):
            if file.endswith('.pptx') or file.endswith('.docx'):
                filepath = os.path.join(storage_dir, file)
                stat = os.stat(filepath)
                templates.append({
                    'name': file,
                    'display_name': file.replace('.pptx', '').replace('.docx', ''),
                    'path': filepath,
                    'size': stat.st_size,
                    'modified': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                    'type': 'PPTX' if file.endswith('.pptx') else 'DOCX',
                    'source': 'stored'
                })
    templates.extend(get_github_templates())
    return templates

def delete_template_file(template_name):
    if os.path.exists(os.path.join(TEMPLATES_DIR, template_name)):
        st.warning("Cannot delete repository templates (documents_template/)")
        return False
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, template_name)
    if os.path.exists(filepath):
        os.remove(filepath)
        config_name = template_name.replace('.pptx', '').replace('.docx', '') + '_config.json'
        config_path = os.path.join(storage_dir, config_name)
        if os.path.exists(config_path):
            os.remove(config_path)
        temp_config = get_temp_config_path(template_name)
        if os.path.exists(temp_config):
            os.remove(temp_config)
        get_saved_templates.clear()
        load_template_from_file.clear()
        return True
    return False

def save_config_to_file(config_data, config_name="template_config.json"):
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, config_name)
    with open(filepath, 'w', encoding='utf-8') as f:
        json.dump(config_data, f, indent=4)
    return filepath

def load_config_from_file(config_name="template_config.json"):
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, config_name)
    if os.path.exists(filepath):
        with open(filepath, 'r', encoding='utf-8') as f:
            return json.load(f)
    return None

def auto_save_config():
    if st.session_state.saved_template_name and st.session_state.custom_mapping:
        config_name = st.session_state.saved_template_name.replace('.pptx', '').replace('.docx', '') + '_config.json'
        save_config_to_file(st.session_state.custom_mapping, config_name)

# --- CTA PRESET FUNCTIONS ---
def detect_cta_sets():
    """Detect CTA sets in the template placeholders"""
    cta_sets = {}
    for token in st.session_state.tokens:
        clean_label = token.replace("{", "").replace("}", "").upper()
        match = re.match(r'CTA(\d+)_(NAME|CONTACT_NUMBER|EMAIL)', clean_label)
        if match:
            cta_num = int(match.group(1))
            field_type = match.group(2)
            if cta_num not in cta_sets:
                cta_sets[cta_num] = {
                    'tokens': {},
                    'fields': set()
                }
            cta_sets[cta_num]['tokens'][field_type] = token
            cta_sets[cta_num]['fields'].add(field_type)
    return cta_sets

def apply_cta_preset_autofill(cta_num, advisor_name):
    """Auto-fill CTA preset values to a specific CTA set"""
    if advisor_name not in contacts_database:
        return False

    contact_info = contacts_database[advisor_name]
    cta_sets = detect_cta_sets()

    if cta_num not in cta_sets:
        return False

    tokens = cta_sets[cta_num]['tokens']

    # Store current values for all tokens before modifying
    current_values = {}
    for token in st.session_state.tokens:
        val_key = f"val_{token}"
        if val_key in st.session_state:
            current_values[token] = st.session_state[val_key]

    # Apply CTA values
    if 'NAME' in tokens:
        st.session_state[f"val_{tokens['NAME']}"] = advisor_name
        st.session_state.temp_form_data[tokens['NAME']] = advisor_name
    if 'CONTACT_NUMBER' in tokens:
        st.session_state[f"val_{tokens['CONTACT_NUMBER']}"] = contact_info["phone"]
        st.session_state.temp_form_data[tokens['CONTACT_NUMBER']] = contact_info["phone"]
    if 'EMAIL' in tokens:
        st.session_state[f"val_{tokens['EMAIL']}"] = contact_info["email"]
        st.session_state.temp_form_data[tokens['EMAIL']] = contact_info["email"]

    # Restore all other values
    cta_tokens = set(tokens.values())
    for token, value in current_values.items():
        val_key = f"val_{token}"
        if token not in cta_tokens:
            if val_key not in st.session_state or st.session_state[val_key] != value:
                st.session_state[val_key] = value
                st.session_state.temp_form_data[token] = value

    if st.session_state.saved_template_name:
        temp_path = get_temp_config_path(st.session_state.saved_template_name)
        try:
            with open(temp_path, 'w', encoding='utf-8') as f:
                json.dump(st.session_state.temp_form_data, f, indent=4)
            return True
        except Exception:
            return False
    return True

# --- BASEMAP CONFIGURATION WITH IMPROVED RELIABILITY ---
BASEMAP_CONFIG = {
    "Satellite (Labels)": {
        "urls": [
            "https://mt1.google.com/vt/lyrs=y&x={x}&y={y}&z={z}",
            "https://mt0.google.com/vt/lyrs=y&x={x}&y={y}&z={z}"
        ],
        "attribution": "Google"
    },
    "Satellite (Streets)": {
        "urls": [
            "https://mt1.google.com/vt/lyrs=y&x={x}&y={y}&z={z}&apistyle=s.t%3A2%7Cp.v%3Aoff",
            "https://mt0.google.com/vt/lyrs=y&x={x}&y={y}&z={z}&apistyle=s.t%3A2%7Cp.v%3Aoff"
        ],
        "attribution": "Google"
    },
    "Satellite (Clean)": {
        "urls": [
            "https://mt1.google.com/vt/lyrs=s&x={x}&y={y}&z={z}",
            "https://mt0.google.com/vt/lyrs=s&x={x}&y={y}&z={z}"
        ],
        "attribution": "Google"
    },
    "Street Map": {
        "urls": [
            "https://mt1.google.com/vt/lyrs=m&x={x}&y={y}&z={z}",
            "https://mt0.google.com/vt/lyrs=m&x={x}&y={y}&z={z}"
        ],
        "attribution": "Google"
    },
    "OSM Carto Light": {
        "urls": [
            "https://a.basemaps.cartocdn.com/light_all/{z}/{x}/{y}.png",
            "https://b.basemaps.cartocdn.com/light_all/{z}/{x}/{y}.png",
            "https://c.basemaps.cartocdn.com/light_all/{z}/{x}/{y}.png"
        ],
        "attribution": "CartoDB"
    },
    "Open Street Map": {
        "urls": [
            "https://a.tile.openstreetmap.org/{z}/{x}/{y}.png",
            "https://b.tile.openstreetmap.org/{z}/{x}/{y}.png",
            "https://c.tile.openstreetmap.org/{z}/{x}/{y}.png"
        ],
        "attribution": "OpenStreetMap"
    }
}

def get_tile_urls(style_name):
    """Get list of tile URLs with failover support"""
    config = BASEMAP_CONFIG.get(style_name)
    if not config:
        return BASEMAP_CONFIG["Street Map"]["urls"]
    return config["urls"]

def get_attribution(style_name):
    """Get attribution for a basemap style"""
    config = BASEMAP_CONFIG.get(style_name)
    return config["attribution"] if config else ""

@st.cache_data(show_spinner=False, ttl=3600, max_entries=4096)
def fetch_tile_with_retry(url_template, zoom, x, y, max_retries=3):
    """Fetch a tile with retry logic and multiple URL fallbacks (cached per tile)."""
    headers = {"User-Agent": "Mozilla/5.0 (compatible; MapGenerator/1.0; +https://example.com)"}
    for attempt in range(max_retries):
        url = url_template.format(z=zoom, x=x, y=y)
        try:
            resp = requests.get(url, headers=headers, timeout=10)
            if resp.status_code == 200:
                return resp.content
            elif resp.status_code == 418:
                continue
        except Exception:
            continue
    return None

# --- DYNAMIC ULTRA HIGH-RESOLUTION BOUNDING BOX GENERATOR ---
def generate_static_map_bounds(n, s, e, w, pin_lat, pin_lon, style="Satellite (Streets)", pin_color="#003366", pin_size=18):
    """Generates high-res map with pin always included - with 1km default radius if no bounds provided"""

    def calculate_1km_bounds(lat, lon):
        """Calculate approximately 1km bounding box centered on the pin"""
        # 1 degree latitude ≈ 111.32 km
        # 1 degree longitude ≈ 111.32 * cos(latitude) km
        lat_deg_per_km = 1.0 / 111.32
        lon_deg_per_km = 1.0 / (111.32 * math.cos(math.radians(lat)))

        # 0.5km in each direction for 1km total (centered on pin)
        lat_offset = lat_deg_per_km * 0.5
        lon_offset = lon_deg_per_km * 0.5

        return lat + lat_offset, lat - lat_offset, lon + lon_offset, lon - lon_offset

    # Check if bounds are valid (not None and have reasonable size)
    bounds_valid = all(x is not None for x in [n, s, e, w])
    if bounds_valid:
        # Check if bounds are too small or invalid
        lat_span = abs(n - s)
        lon_span = abs(e - w)
        # If bounds are extremely small (less than 10 meters), use 1km
        if lat_span < 0.0001 or lon_span < 0.0001:
            bounds_valid = False

    if not bounds_valid:
        # Use 1km radius centered on pin
        n, s, e, w = calculate_1km_bounds(pin_lat, pin_lon)

    lon_span = e - w
    lat_span = n - s
    target_width_tiles = 8
    if lon_span <= 0: lon_span = 0.001
    zoom = int(math.log2((360.0 / lon_span) * target_width_tiles))
    zoom = max(13, min(20, zoom))
    if lon_span < 0.01 and lat_span < 0.01:
        zoom = min(20, zoom + 2)

    def deg2num(lat_deg, lon_deg, z):
        lat_rad = math.radians(lat_deg)
        n_tiles = 2.0 ** z
        xtile = int((lon_deg + 180.0) / 360.0 * n_tiles)
        ytile = int((1.0 - math.log(math.tan(lat_rad) + (1 / math.cos(lat_rad))) / math.pi) / 2.0 * n_tiles)
        return (xtile, ytile)

    x_min, y_min = deg2num(n, w, zoom)
    x_max, y_max = deg2num(s, e, zoom)
    if x_max == x_min: x_max += 1
    if y_max == y_min: y_max += 1
    if (x_max - x_min + 1) * (y_max - y_min + 1) > 100:
        zoom -= 1
        x_min, y_min = deg2num(n, w, zoom)
        x_max, y_max = deg2num(s, e, zoom)
    width_tiles = x_max - x_min + 1
    height_tiles = y_max - y_min + 1
    tile_size = 256
    scale_factor = 2
    stitched = Image.new('RGB', (width_tiles * tile_size * scale_factor, height_tiles * tile_size * scale_factor))
    headers = {"User-Agent": "Mozilla/5.0 (compatible; MapGenerator/1.0; +https://example.com)"}

    tile_urls = get_tile_urls(style)

    # Fetch and stitch tiles
    for x in range(x_min, x_max + 1):
        for y in range(y_min, y_max + 1):
            tile_data = None
            for url_template in tile_urls:
                tile_data = fetch_tile_with_retry(url_template, zoom, x, y)
                if tile_data is not None:
                    break

            if tile_data is not None:
                try:
                    img = Image.open(io.BytesIO(tile_data))
                    img = img.resize((tile_size * scale_factor, tile_size * scale_factor), Image.Resampling.LANCZOS)
                    stitched.paste(img, ((x - x_min) * tile_size * scale_factor, (y - y_min) * tile_size * scale_factor))
                except Exception:
                    pass

    def num2px(lat_deg, lon_deg, z):
        lat_rad = math.radians(lat_deg)
        n_tiles = 2.0 ** z
        px_x = (lon_deg + 180.0) / 360.0 * n_tiles * tile_size * scale_factor
        px_y = (1.0 - math.log(math.tan(lat_rad) + (1 / math.cos(lat_rad))) / math.pi) / 2.0 * n_tiles * tile_size * scale_factor
        return px_x, px_y

    # Calculate pixel positions for cropping
    px_w, py_n = num2px(n, w, zoom)
    px_e, py_s = num2px(s, e, zoom)
    base_x = x_min * tile_size * scale_factor
    base_y = y_min * tile_size * scale_factor

    # For no bounds case, ensure the crop is perfectly centered on pin
    if not bounds_valid:
        # Calculate pin pixel position in the stitched image
        pin_px_x, pin_px_y = num2px(pin_lat, pin_lon, zoom)

        # Calculate crop dimensions
        crop_width = int(px_e - px_w)
        crop_height = int(py_s - py_n)

        # Center the crop on the pin
        left = int(pin_px_x - base_x - crop_width // 2)
        top = int(pin_px_y - base_y - crop_height // 2)
        right = left + crop_width
        bottom = top + crop_height

        # Ensure we don't go out of bounds
        stitched_width = width_tiles * tile_size * scale_factor
        stitched_height = height_tiles * tile_size * scale_factor

        if left < 0:
            right -= left
            left = 0
        if top < 0:
            bottom -= top
            top = 0
        if right > stitched_width:
            left -= (right - stitched_width)
            right = stitched_width
        if bottom > stitched_height:
            top -= (bottom - stitched_height)
            bottom = stitched_height

        # Ensure valid crop dimensions
        if right <= left: right = left + 100
        if bottom <= top: bottom = top + 100

        cropped = stitched.crop((left, top, right, bottom)).convert("RGBA")
    else:
        # Use the bounds-based crop
        left = int(px_w - base_x)
        top = int(py_n - base_y)
        right = int(px_e - base_x)
        bottom = int(py_s - base_y)
        if right <= left: right = left + 100
        if bottom <= top: bottom = top + 100
        cropped = stitched.crop((left, top, right, bottom)).convert("RGBA")

    # --- DRAW PIN MARKER (ALWAYS INCLUDED AND CENTERED FOR NO BOUNDS) ---
    draw = ImageDraw.Draw(cropped)
    pin_px_x, pin_px_y = num2px(pin_lat, pin_lon, zoom)

    # Calculate pin position relative to the crop
    pin_local_x = int(pin_px_x - base_x) - left
    pin_local_y = int(pin_px_y - base_y) - top

    # For no bounds case, the pin should be exactly in the center
    if not bounds_valid:
        pin_local_x = cropped.width // 2
        pin_local_y = cropped.height // 2

    # Ensure pin is within bounds (clamp if needed)
    pin_local_x = max(0, min(pin_local_x, cropped.width - 1))
    pin_local_y = max(0, min(pin_local_y, cropped.height - 1))

    # Draw pin with shadow and star
    radius = int((pin_size / 2) * scale_factor)
    shadow_offset = max(1, int(radius * 0.15))

    # Shadow
    draw.ellipse([
        pin_local_x - radius - shadow_offset,
        pin_local_y - radius - shadow_offset,
        pin_local_x + radius + shadow_offset,
        pin_local_y + radius + shadow_offset
    ], fill=(0, 0, 0, 60))

    # Pin circle
    draw.ellipse([
        pin_local_x - radius,
        pin_local_y - radius,
        pin_local_x + radius,
        pin_local_y + radius
    ], fill=pin_color, outline=(255, 255, 255), width=max(1, int(radius * 0.1)))

    # Star
    star_size = int(radius * 0.55)
    star_points = []
    for i in range(10):
        angle = math.radians(i * 36 - 90)
        r = star_size if i % 2 == 0 else star_size * 0.4
        star_points.append((pin_local_x + r * math.cos(angle), pin_local_y + r * math.sin(angle)))
    draw.polygon(star_points, fill=(255, 255, 255))

    # Add a small glow effect around the pin for visibility
    glow_radius = int(radius * 1.5)
    for i in range(3):
        alpha = 30 - i * 10
        glow_radius_i = glow_radius + i * 5
        draw.ellipse([
            pin_local_x - glow_radius_i,
            pin_local_y - glow_radius_i,
            pin_local_x + glow_radius_i,
            pin_local_y + glow_radius_i
        ], outline=(255, 255, 255, alpha), width=1)

    final_img = cropped.convert("RGB")
    img_byte_arr = io.BytesIO()
    final_img.save(img_byte_arr, format='PNG', quality=100, optimize=True)
    img_byte_arr.seek(0)
    return img_byte_arr

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# --- ISOLATED FULL-SCREEN MAP EDITOR PAGE ---
def render_isolated_map_editor():
    token_key = st.session_state.active_map_editor_token

    st.markdown("""
        <style>
            div[data-testid="stHorizontalBlock"] {
                align-items: flex-end !important;
                gap: 12px !important;
            }
            div[data-baseweb="input"], div[data-baseweb="select"], .stColorPicker div {
                height: 38px !important;
            }
            .manual-picker-label {
                font-family: 'Montserrat', sans-serif !important;
                font-size: 14px !important;
                color: #003366 !important;
                margin-bottom: 8px !important;
                line-height: 1.2;
            }
        </style>
    """, unsafe_allow_html=True)

    st.markdown('<div class="editor-card">', unsafe_allow_html=True)
    col_back, col_title = st.columns([1, 4])
    with col_back:
        def return_to_main():
            if st.session_state.saved_template_name:
                temp_path = get_temp_config_path(st.session_state.saved_template_name)
                if os.path.exists(temp_path):
                    try:
                        with open(temp_path, 'r', encoding='utf-8') as f:
                            loaded_data = json.load(f)
                            st.session_state.temp_form_data = loaded_data
                            for token, value in loaded_data.items():
                                current_type = st.session_state.custom_mapping.get(token, "Text")
                                if current_type != "Image":
                                    st.session_state[f"val_{token}"] = value
                    except Exception as e:
                        st.error(f"Error restoring data: {e}")

            st.session_state.restore_form_data = True
            st.session_state.active_map_editor_token = None
            st.rerun()

        if st.button("Back to Document", key="back_from_map", on_click=return_to_main):
            pass

    with col_title:
        st.markdown(f"### Map Editor: {token_key}")
    st.markdown("</div><br>", unsafe_allow_html=True)

    style_key = f"map_style_{token_key}"
    coord_key = f"map_coord_{token_key}"
    color_key = f"map_color_{token_key}"
    size_key = f"map_size_{token_key}"
    dragged_key = f"map_dragged_{token_key}"
    image_key = f"map_bytes_holder_{token_key}"
    bounds_key = f"map_bounds_{token_key}"
    export_trigger_key = f"map_export_active_{token_key}"

    if style_key not in st.session_state: st.session_state[style_key] = "Satellite (Streets)"
    if coord_key not in st.session_state: st.session_state[coord_key] = "14.5995, 120.9842"
    if color_key not in st.session_state: st.session_state[color_key] = "#003366"
    if size_key not in st.session_state: st.session_state[size_key] = 20
    if image_key not in st.session_state: st.session_state[image_key] = None
    if bounds_key not in st.session_state: st.session_state[bounds_key] = None
    if export_trigger_key not in st.session_state: st.session_state[export_trigger_key] = False

    if dragged_key in st.session_state:
        st.session_state[coord_key] = st.session_state[dragged_key]
        del st.session_state[dragged_key]

    map_styles = ["Satellite (Streets)", "Satellite (Labels + Streets)", "Satellite (Clean)",
                  "Street Map", "OSM Carto Light", "Open Street Map"]

    c_btn, c_style, c_color, c_size, c_coord = st.columns([1.4, 1.8, 0.8, 1.0, 2.8])
    with c_btn:
        st.markdown("<div style='margin-bottom: 2px;'></div>", unsafe_allow_html=True)
        export_clicked = st.button("Confirm and Export", type="primary", key=f"export_map_{token_key}", use_container_width=True)
        if export_clicked:
            st.session_state[export_trigger_key] = True

    with c_style:
        basemap_style = st.selectbox(label="Basemap Layer", options=map_styles, key=style_key)
    with c_color:
        st.markdown('<div class="manual-picker-label">Pin Color</div>', unsafe_allow_html=True)
        pin_color = st.color_picker(label="Pin Color", key=color_key, label_visibility="collapsed")
    with c_size:
        st.markdown('<div class="manual-picker-label">Pin Size</div>', unsafe_allow_html=True)
        pin_size = st.number_input(label="Pin Size", min_value=8, max_value=64, step=1, value=st.session_state[size_key], key=size_key, label_visibility="collapsed")
    with c_coord:
        coord_input = st.text_input(label="Enter Coordinates", key=coord_key, placeholder="Lat, Lon")

    st.markdown("<div style='margin-bottom: 16px;'></div>", unsafe_allow_html=True)
    try:
        plat, plon = map(float, coord_input.split(","))
    except ValueError:
        plat, plon = 14.5995, 120.9842

    if st.session_state[export_trigger_key]:
        with st.spinner("Exporting Map... Please wait"):
            n, s, e, w = None, None, None, None

            # Check if user drew a rectangle
            if st.session_state.get(bounds_key):
                b = st.session_state[bounds_key]
                if b and "_northEast" in b and "_southWest" in b:
                    n, s = b["_northEast"]["lat"], b["_southWest"]["lat"]
                    e, w = b["_northEast"]["lng"], b["_southWest"]["lng"]

                    # Validate bounds have reasonable size
                    if abs(n - s) < 0.0001 or abs(e - w) < 0.0001:
                        n, s, e, w = None, None, None, None

            # If no valid bounds, use 1km radius (handled in generate_static_map_bounds)
            # Pass None values - the function will calculate 1km bounds

            map_img_bytes = generate_static_map_bounds(
                n=n, s=s, e=e, w=w,
                pin_lat=plat, pin_lon=plon,
                style=basemap_style,
                pin_color=pin_color,
                pin_size=int(pin_size)
            )

            # Store the generated map image
            st.session_state[image_key] = map_img_bytes
            st.session_state[f"coord_{token_key}"] = f"{plat}, {plon}"

            # Save to temp data
            if st.session_state.temp_form_data:
                st.session_state.temp_form_data[token_key] = f"{plat}, {plon}"
                temp_path = get_temp_config_path(st.session_state.saved_template_name)
                with open(temp_path, 'w', encoding='utf-8') as f:
                    json.dump(st.session_state.temp_form_data, f, indent=4)

            st.session_state[export_trigger_key] = False
            st.session_state.restore_form_data = True
            st.session_state.active_map_editor_token = None
            st.success(f"Map with pin attached successfully!")
            time.sleep(0.5)
            st.rerun()

    tiles_dict = {}
    attr_dict = {}
    for style in map_styles:
        urls = get_tile_urls(style)
        tiles_dict[style] = urls[0] if urls else ""
        attr_dict[style] = get_attribution(style)

    m = folium.Map(location=[plat, plon], zoom_start=15, tiles=tiles_dict[basemap_style], attr=attr_dict[basemap_style], zoom_control=True)

    icon_html = f"""
    <div style="position: relative; width: {pin_size}px; height: {pin_size}px;">
        <svg width="{pin_size}" height="{pin_size}" viewBox="0 0 40 40" style="width: 100%; height: 100%;">
            <defs><filter id="shadow" x="-20%" y="-20%" width="140%" height="140%"><feDropShadow dx="0" dy="2" stdDeviation="2" flood-opacity="0.3"/></filter></defs>
            <g filter="url(#shadow)">
                <circle cx="20" cy="20" r="18" fill="{pin_color}" stroke="white" stroke-width="2"/>
                <polygon points="20,6 23.5,14.5 32,15 25.5,21 27.5,30 20,25 12.5,30 14.5,21 8,15 16.5,14.5" fill="white"/>
            </g>
        </svg>
    </div>
    """
    folium.Marker([plat, plon], draggable=True, icon=folium.DivIcon(html=icon_html)).add_to(m)

    Draw(
        export=False,
        position='topleft',
        draw_options={
            'polyline': False,
            'polygon': False,
            'circle': False,
            'marker': False,
            'circlemarker': False,
            'rectangle': True
        },
        edit_options={'edit': True}
    ).add_to(m)

    st.info("Draw a rectangle to define crop area (visible for guidance only, not in export) | Drag pin to reposition")
    map_data = st_folium(
        m, height=600, width=1300, use_container_width=True, key=f"int_map_{token_key}",
        returned_objects=["last_active_drawing", "bounds", "last_marker_moved"]
    )

    if isinstance(map_data, dict):
        if map_data.get("bounds"):
            st.session_state[bounds_key] = map_data["bounds"]
        if map_data.get("last_marker_moved"):
            moved = map_data["last_marker_moved"]
            if moved:
                new_coord = f"{round(moved['lat'], 5)}, {round(moved['lng'], 5)}"
                if new_coord != st.session_state.get(coord_key, ""):
                    st.session_state[dragged_key] = new_coord
                    st.rerun()

# --- CORE UTILITIES ---
def smart_crop_to_fit(img_file, target_w_emu, target_h_emu):
    try:
        img = Image.open(img_file)
        img_w, img_h = img.size
        target_ratio = target_w_emu / target_h_emu
        if img_w / img_h > target_ratio:
            new_w = int(img_h * target_ratio)
            left = (img_w - new_w) // 2
            img = img.crop((left, 0, left + new_w, img_h))
        else:
            new_h = int(img_w / target_ratio)
            top = (img_h - new_h) // 2
            img = img.crop((0, top, img_w, top + new_h))
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='PNG', quality=95)
        img_byte_arr.seek(0)
        return img_byte_arr
    except Exception:
        return img_file

def extract_placeholders_from_pptx(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    tokens, seen = [], set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for token in re.findall(r'\{\{.*?\}\}', shape.text):
                    if token not in seen: tokens.append(token); seen.add(token)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for token in re.findall(r'\{\{.*?\}\}', cell.text):
                            if token not in seen: tokens.append(token); seen.add(token)
    return tokens

def extract_placeholders_from_docx(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    tokens, seen = [], set()
    for paragraph in doc.paragraphs:
        for token in re.findall(r'\{\{.*?\}\}', paragraph.text):
            if token not in seen: tokens.append(token); seen.add(token)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for token in re.findall(r'\{\{.*?\}\}', cell.text):
                    if token not in seen: tokens.append(token); seen.add(token)
    return tokens

@st.cache_data(show_spinner=False, ttl=600)
def extract_placeholders(template_bytes, template_type):
    if template_type == 'pptx': return extract_placeholders_from_pptx(template_bytes)
    if template_type == 'docx': return extract_placeholders_from_docx(template_bytes)
    return []

def clean_empty_placeholders(text):
    """Remove any {{...}} placeholders that remain empty and clean up whitespace"""
    if not text:
        return text
    cleaned = re.sub(r'\{\{[^}]*\}\}', '', text)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    return cleaned

def replace_text_in_paragraph(paragraph, text_inputs):
    """Replace text in paragraph, removing empty placeholders"""
    # Process each run
    for run in paragraph.runs:
        current_text = run.text
        for token, value in text_inputs.items():
            if token in current_text:
                if value and str(value).strip():  # Value exists and is not empty
                    current_text = current_text.replace(token, str(value))
                else:  # Value is empty, remove the placeholder
                    current_text = current_text.replace(token, '')
        # Clean up any remaining placeholders and extra whitespace
        run.text = clean_empty_placeholders(current_text)

    # Handle paragraph text if no runs exist or if there are still placeholders
    if hasattr(paragraph, 'text') and paragraph.text:
        current_text = paragraph.text
        # Check if any placeholders remain
        has_placeholder = any(token in current_text for token in text_inputs.keys())
        if has_placeholder:
            for token, value in text_inputs.items():
                if token in current_text:
                    if value and str(value).strip():
                        current_text = current_text.replace(token, str(value))
                    else:
                        current_text = current_text.replace(token, '')
            current_text = clean_empty_placeholders(current_text)

            if not paragraph.runs:
                paragraph.add_run(current_text)
            else:
                # Update the first run with cleaned text
                for run in paragraph.runs:
                    if any(token in run.text for token in text_inputs.keys()):
                        run.text = current_text
                        break

def generate_pptx_bytes(template_bytes, text_inputs, image_inputs):
    prs = Presentation(io.BytesIO(template_bytes))
    for slide in prs.slides:
        shapes_to_delete, images_to_add = [], []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for img_token, img_file in image_inputs.items():
                    if img_token in shape.text and img_file is not None:
                        images_to_add.append((img_file, shape.left, shape.top, shape.width, shape.height))
                        shapes_to_delete.append(shape)
                        break
        for shape in slide.shapes:
            if shape not in shapes_to_delete:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        replace_text_in_paragraph(paragraph, text_inputs)
                if hasattr(shape, 'table') and shape.table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    replace_text_in_paragraph(paragraph, text_inputs)
        for img_file, left, top, width, height in images_to_add:
            try:
                slide.shapes.add_picture(smart_crop_to_fit(img_file, width, height), left, top, width=width, height=height)
            except Exception:
                pass
        for old_shape in shapes_to_delete:
            try:
                sp = old_shape._element
                sp.getparent().remove(sp)
            except Exception:
                pass
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    return pptx_stream.getvalue()

def generate_docx_bytes(template_bytes, text_inputs, image_inputs):
    doc = Document(io.BytesIO(template_bytes))
    for paragraph in doc.paragraphs:
        if not any(img_token in paragraph.text for img_token in image_inputs.keys()):
            replace_text_in_paragraph(paragraph, text_inputs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.paragraphs:
                    for paragraph in cell.paragraphs:
                        replace_text_in_paragraph(paragraph, text_inputs)
    doc_stream = io.BytesIO()
    doc.save(doc_stream)
    return doc_stream.getvalue()


def docx_bytes_to_pdf(docx_bytes):
    """Convert generated DOCX bytes to PDF bytes.

    Engine auto-detect: try `docx2pdf` (wraps installed Microsoft Word / LibreOffice)
    first, then fall back to a LibreOffice `soffice` headless subprocess. Raises a
    clear error if no usable converter is installed.
    """
    import shutil
    import subprocess
    import tempfile

    with tempfile.TemporaryDirectory() as td:
        docx_path = os.path.join(td, "input.docx")
        pdf_path = os.path.join(td, "output.pdf")
        with open(docx_path, "wb") as f:
            f.write(docx_bytes)

        # Engine 1: docx2pdf (uses installed MS Word, or LibreOffice if configured).
        try:
            import docx2pdf
        except Exception:
            docx2pdf = None

        if docx2pdf is not None:
            try:
                docx2pdf.convert(docx_path, pdf_path)
                if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                    with open(pdf_path, "rb") as f:
                        return f.read()
            except Exception:
                pass  # fall through to LibreOffice

        # Engine 2: LibreOffice headless.
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice is not None:
            try:
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf", "--outdir", td, docx_path],
                    capture_output=True,
                    timeout=180,
                )
                # LibreOffice writes output.pdf in `td`.
                if result.returncode == 0 and os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                    with open(pdf_path, "rb") as f:
                        return f.read()
            except Exception:
                pass

        raise RuntimeError(
            "No DOCX→PDF converter available. Install Microsoft Word (docx2pdf) or "
            "LibreOffice (`soffice`) on this machine to export PDFs."
        )


def get_download_filename(template_name, file_type):
    """Generate filename: Generated_TemplateName_Date"""
    # Remove template_ prefix and file extension
    base_name = re.sub(r'^template_', '', template_name or "Document")
    base_name = re.sub(r'\.(pptx|docx)$', '', base_name)
    base_name = re.sub(r'[^\w\-_. ]', '_', base_name)

    # Get current date in MMDDYYYY format
    current_date = datetime.now().strftime('%m%d%Y')

    # Format: Generated_TemplateName_Date.ext
    return f"Generated_{base_name}_{current_date}.{file_type}"

def autosave_current_form_data():
    if not st.session_state.saved_template_name or not st.session_state.tokens:
        return False
    for token in st.session_state.tokens:
        val_key = f"val_{token}"
        if val_key in st.session_state:
            current_type = st.session_state.custom_mapping.get(token, "Text")
            if current_type != "Image":
                st.session_state.temp_form_data[token] = st.session_state[val_key]
    temp_path = get_temp_config_path(st.session_state.saved_template_name)
    try:
        with open(temp_path, 'w', encoding='utf-8') as f:
            json.dump(st.session_state.temp_form_data, f, indent=4)
        return True
    except Exception:
        return False

def restore_form_data_from_session():
    if not st.session_state.saved_template_name:
        return False
    has_session_data = any(f"val_{token}" in st.session_state for token in st.session_state.tokens)
    if has_session_data:
        return True
    if st.session_state.temp_form_data:
        for token, value in st.session_state.temp_form_data.items():
            current_type = st.session_state.custom_mapping.get(token, "Text")
            if current_type != "Image" and f"val_{token}" not in st.session_state:
                st.session_state[f"val_{token}"] = value
        return True
    temp_path = get_temp_config_path(st.session_state.saved_template_name)
    if os.path.exists(temp_path):
        try:
            with open(temp_path, 'r', encoding='utf-8') as f:
                loaded_data = json.load(f)
                st.session_state.temp_form_data = loaded_data
                for token, value in loaded_data.items():
                    current_type = st.session_state.custom_mapping.get(token, "Text")
                    if current_type != "Image" and f"val_{token}" not in st.session_state:
                        st.session_state[f"val_{token}"] = value
                return True
        except Exception:
            pass
    return False

def purge_all_temporary_data():
    if st.session_state.saved_template_name:
        temp_path = get_temp_config_path(st.session_state.saved_template_name)
        if os.path.exists(temp_path):
            try: os.remove(temp_path)
            except Exception: pass
    if st.session_state.tokens:
        for token in st.session_state.tokens:
            if f"val_{token}" in st.session_state: del st.session_state[f"val_{token}"]
            if f"map_bytes_holder_{token}" in st.session_state: del st.session_state[f"map_bytes_holder_{token}"]
    st.session_state.temp_form_data = {}

# --- STATE INIT ---
if "active_map_editor_token" not in st.session_state: st.session_state.active_map_editor_token = None
if "custom_mapping" not in st.session_state: st.session_state.custom_mapping = {}
if "tokens" not in st.session_state: st.session_state.tokens = []
if "template_bytes" not in st.session_state: st.session_state.template_bytes = None
if "saved_template_name" not in st.session_state: st.session_state.saved_template_name = None
if "template_loaded" not in st.session_state: st.session_state.template_loaded = False
if "template_type" not in st.session_state: st.session_state.template_type = None
if "show_delete_confirm" not in st.session_state: st.session_state.show_delete_confirm = False
if "template_to_delete" not in st.session_state: st.session_state.template_to_delete = None
if "save_success" not in st.session_state: st.session_state.save_success = False
if "saved_file_name" not in st.session_state: st.session_state.saved_file_name = None
if "clear_uploader" not in st.session_state: st.session_state.clear_uploader = False
if "restore_form_data" not in st.session_state: st.session_state.restore_form_data = False
if "show_type_mapping" not in st.session_state: st.session_state.show_type_mapping = False
if "temp_form_data" not in st.session_state: st.session_state.temp_form_data = {}

# --- PAGE HEADER (native) ---
st.markdown('<p class="page-eyebrow">Documents</p>', unsafe_allow_html=True)
st.markdown('<p class="docs-title">Generator</p>', unsafe_allow_html=True)
st.markdown(
    '<p class="docs-caption">Generate branded documents from PPTX/DOCX templates with text, image, and map placeholders.</p>',
    unsafe_allow_html=True,
)
st.markdown(
    '<div class="local-only-note">Documents are generated and exported locally — nothing is saved to the database.</div>',
    unsafe_allow_html=True,
)

# --- APP ROUTER ---
if st.session_state.active_map_editor_token:
    render_isolated_map_editor()
else:
    if st.session_state.restore_form_data:
        restore_form_data_from_session()
        st.session_state.restore_form_data = False

    st.markdown("<hr style='margin: 4px 0 12px 0;'>", unsafe_allow_html=True)
    st.markdown('<div class="section-header">Templates</div>', unsafe_allow_html=True)

    col_template1, col_template2 = st.columns(2)
    with col_template1:
        saved_templates = get_saved_templates()
        template_options = ["Select saved template"]
        github_templates = [t for t in saved_templates if t['source'] == 'github']
        stored_templates = [t for t in saved_templates if t['source'] == 'stored']

        if github_templates:
            template_options.append("--- Templates ---")
            for t in github_templates: template_options.append(f"{t['display_name']} ({t['type']})")
        if stored_templates:
            template_options.append("--- User Uploaded Templates ---")
            for t in stored_templates: template_options.append(f"{t['display_name']} ({t['type']})")

        dropdown_col, delete_col = st.columns([4, 1])
        with dropdown_col:
            selected_template = st.selectbox("Load Template", template_options, key="saved_template_select", label_visibility="collapsed")
        with delete_col:
            if selected_template and selected_template != "Select saved template" and not selected_template.startswith("---"):
                template_display = selected_template.split(' (')[0].strip()
                for t in saved_templates:
                    if t['display_name'] == template_display:
                        if t['source'] == 'stored' and st.button("Delete", key="delete_template"):
                            st.session_state.show_delete_confirm = True
                            st.session_state.template_to_delete = t['name']
                            st.rerun()
                        break

        if st.session_state.show_delete_confirm:
            st.warning(f"Are you sure you want to delete '{st.session_state.template_to_delete}'?")
            col_confirm1, col_confirm2 = st.columns(2)
            with col_confirm1:
                if st.button("Yes, Delete", key="confirm_delete"):
                    if delete_template_file(st.session_state.template_to_delete):
                        st.session_state.template_bytes = None
                        st.session_state.saved_template_name = None
                        st.session_state.template_loaded = False
                        st.session_state.tokens = []
                        st.session_state.temp_form_data = {}
                        st.session_state.show_delete_confirm = False
                        st.session_state.template_to_delete = None
                        st.rerun()
            with col_confirm2:
                if st.button("Cancel", key="cancel_delete"):
                    st.session_state.show_delete_confirm = False
                    st.session_state.template_to_delete = None
                    st.rerun()

        if selected_template and selected_template != "Select saved template" and not selected_template.startswith("---"):
            template_display = selected_template.split(' (')[0].strip()
            for t in saved_templates:
                if t['display_name'] == template_display:
                    template_name = t['name']
                    template_bytes = load_template_from_file(template_name)
                    if template_bytes:
                        if st.session_state.saved_template_name != template_name:
                            st.session_state.temp_form_data = {}
                        st.session_state.template_bytes = template_bytes
                        st.session_state.saved_template_name = template_name
                        st.session_state.template_loaded = True
                        st.session_state.template_type = 'pptx' if template_name.endswith('.pptx') else 'docx'
                        config_data = load_config_from_file(template_name.replace('.pptx', '').replace('.docx', '') + '_config.json')
                        if config_data: st.session_state.custom_mapping = config_data
                        st.session_state.tokens = extract_placeholders(template_bytes, st.session_state.template_type)
                        restore_form_data_from_session()
                    break

    with col_template2:
        uploader_key = "new_template_upload_clear" if st.session_state.clear_uploader else "new_template_upload"
        uploaded_template = st.file_uploader("Upload New Template", type=["pptx", "docx"], label_visibility="collapsed", key=uploader_key)
        if st.session_state.clear_uploader: st.session_state.clear_uploader = False

        if uploaded_template:
            template_bytes = uploaded_template.getvalue()
            st.session_state.template_bytes = template_bytes
            st.session_state.saved_template_name = uploaded_template.name
            st.session_state.template_loaded = True
            st.session_state.template_type = 'pptx' if uploaded_template.name.endswith('.pptx') else 'docx'
            st.session_state.tokens = extract_placeholders(template_bytes, st.session_state.template_type)
            st.session_state.temp_form_data = {}

            if st.button("Save Template", key="save_template_btn", use_container_width=True):
                save_template_to_file(template_bytes, uploaded_template.name)
                if st.session_state.custom_mapping:
                    save_config_to_file(st.session_state.custom_mapping, uploaded_template.name.replace('.pptx', '').replace('.docx', '') + '_config.json')
                st.session_state.save_success = True
                st.session_state.saved_file_name = uploaded_template.name
                st.session_state.clear_uploader = True
                st.rerun()

    if st.session_state.save_success:
        st.success(f"Template '{st.session_state.saved_file_name}' saved successfully!")
        st.session_state.save_success = False

    if st.session_state.template_bytes is not None:
        template_name = st.session_state.saved_template_name or "Unsaved Template"
        is_github = os.path.exists(os.path.join(TEMPLATES_DIR, template_name))
        st.markdown(f'<div class="saved-indicator">Active: {template_name}{"" if is_github else ""} ({st.session_state.template_type.upper()})</div>', unsafe_allow_html=True)

    text_data, image_data, field_types = {}, {}, {}

    if st.session_state.template_bytes is not None and st.session_state.tokens:
        tokens = st.session_state.tokens

        # --- CTA PRESETS - SIMPLIFIED ---
        cta_sets = detect_cta_sets()
        if cta_sets:
            st.markdown("**CTA Presets**")
            num_ctas = len(cta_sets)
            cols = st.columns(num_ctas)

            for idx, cta_num in enumerate(sorted(cta_sets.keys())):
                with cols[idx]:
                    cta_name_token = cta_sets[cta_num]['tokens'].get('NAME')
                    current_advisor = ""
                    if cta_name_token and f"val_{cta_name_token}" in st.session_state:
                        current_advisor = st.session_state[f"val_{cta_name_token}"]

                    st.caption(f"CTA{cta_num}")

                    selected_advisor = st.selectbox(
                        f"Select advisor",
                        options=[""] + list(contacts_database.keys()),
                        index=list(contacts_database.keys()).index(current_advisor) + 1 if current_advisor in contacts_database else 0,
                        key=f"cta_autofill_{cta_num}",
                        label_visibility="collapsed"
                    )

                    if selected_advisor and selected_advisor != current_advisor:
                        apply_cta_preset_autofill(cta_num, selected_advisor)
                        st.rerun()
            st.markdown("---")

        with st.expander("Data Type Mapping", expanded=st.session_state.show_type_mapping):
            st.markdown("Configure the data type for each placeholder field.")

            # Simple 2-column layout for type mapping
            cols = st.columns(2)
            for idx, token in enumerate(tokens):
                with cols[idx % 2]:
                    clean_label = token.replace("{", "").replace("}", "")
                    current_type = st.session_state.custom_mapping.get(token, "Text")

                    col_lbl, col_sel = st.columns([1, 1.5])
                    with col_lbl:
                        st.markdown(f'<span style="font-size:12px; font-weight:500;">{clean_label}</span>', unsafe_allow_html=True)
                    with col_sel:
                        data_type = st.selectbox(
                            "",
                            ["Text", "Image", "Map"],
                            index=["Text", "Image", "Map"].index(current_type) if current_type in ["Text", "Image", "Map"] else 0,
                            key=f"type_mapping_{token}",
                            label_visibility="collapsed"
                        )
                        if data_type != current_type:
                            st.session_state.custom_mapping[token] = data_type
                            auto_save_config()
                            st.rerun()

        st.markdown('<div class="section-header">Placeholder Values</div>', unsafe_allow_html=True)

        # --- RENDER FIELDS IN ORIGINAL ORDER ---
        for idx, token in enumerate(tokens):
            col_target = idx % 2
            if col_target == 0:
                ui_col_1, ui_col_2 = st.columns(2)
                current_block_column = ui_col_1
            else:
                current_block_column = ui_col_2

            with current_block_column:
                clean_label = token.replace("{", "").replace("}", "")
                current_type = st.session_state.custom_mapping.get(token, "Text")

                clean_upper = clean_label.upper()
                cta_match = re.match(r'CTA(\d+)_(NAME|CONTACT_NUMBER|EMAIL)', clean_upper)
                is_cta = cta_match is not None

                label_text = clean_label
                if is_cta:
                    cta_num = cta_match.group(1)
                    label_text = f"{clean_label} (CTA{cta_num})"

                st.markdown(f'<div class="placeholder-label">{label_text}</div>', unsafe_allow_html=True)

                if current_type == "Image" and st.session_state.template_type == 'pptx':
                    image_data[token] = st.file_uploader(clean_label, type=["png", "jpg", "jpeg"], key=f"val_{token}", label_visibility="collapsed")
                    field_types[token] = "Image"

                elif current_type == "Map" and st.session_state.template_type == 'pptx':
                    saved_map_img = st.session_state.get(f"map_bytes_holder_{token}")
                    if saved_map_img:
                        image_data[token] = saved_map_img
                        st.caption("Map attached.")

                    def save_all_and_navigate(token_key):
                        for t in st.session_state.tokens:
                            val_key = f"val_{t}"
                            if val_key in st.session_state:
                                current_type_check = st.session_state.custom_mapping.get(t, "Text")
                                if current_type_check != "Image":
                                    st.session_state.temp_form_data[t] = st.session_state[val_key]

                        if st.session_state.saved_template_name:
                            temp_path = get_temp_config_path(st.session_state.saved_template_name)
                            try:
                                with open(temp_path, 'w', encoding='utf-8') as f:
                                    json.dump(st.session_state.temp_form_data, f, indent=4)
                            except Exception: pass

                        st.session_state.active_map_editor_token = token_key
                        st.rerun()

                    if st.button("Open Map Editor", key=f"btn_map_{token}", use_container_width=True, on_click=save_all_and_navigate, args=(token,)):
                        pass

                    field_types[token] = "Image"

                else:
                    if current_type in ["Image", "Map"] and st.session_state.template_type != 'pptx':
                        st.warning("Images/Maps are only supported in PPTX files.")

                    current_value = ""
                    if f"val_{token}" in st.session_state:
                        current_value = st.session_state[f"val_{token}"]
                    elif token in st.session_state.temp_form_data:
                        current_value = st.session_state.temp_form_data[token]
                        st.session_state[f"val_{token}"] = current_value

                    new_value = st.text_input(
                        "",
                        value=current_value,
                        key=f"val_{token}",
                        label_visibility="collapsed",
                        placeholder=f"Enter {clean_label}..."
                    )

                    if new_value != current_value:
                        st.session_state[f"val_{token}"] = new_value
                        st.session_state.temp_form_data[token] = new_value

                        if st.session_state.saved_template_name:
                            temp_path = get_temp_config_path(st.session_state.saved_template_name)
                            try:
                                with open(temp_path, 'w', encoding='utf-8') as f:
                                    json.dump(st.session_state.temp_form_data, f, indent=4)
                            except Exception: pass

                    text_data[token] = new_value
                    field_types[token] = "Text"

                st.markdown('<div style="margin-bottom:14px;"></div>', unsafe_allow_html=True)

    # --- DOWNLOAD & CLEANUP SECTION ---
    if st.session_state.template_bytes is not None:
        st.markdown('<div class="section-header">Download Document</div>', unsafe_allow_html=True)

        base_template_name = re.sub(r'\.(pptx|docx)$', '', st.session_state.saved_template_name or "Generated_Document")
        col1, col2 = st.columns(2)

        with col1:
            if st.session_state.template_type != 'pptx':
                st.button("Download PPTX", disabled=True, use_container_width=True)
            else:
                try:
                    pptx_data = generate_pptx_bytes(st.session_state.template_bytes, text_data, image_data)
                    st.download_button(
                        label="Download PPTX", data=pptx_data, file_name=get_download_filename(base_template_name, "pptx"),
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True, key="download_pptx", on_click=purge_all_temporary_data
                    )
                except Exception as e:
                    st.error(f"Error generating PPTX: {str(e)}")

        with col2:
            if st.session_state.template_type != 'docx':
                st.button("Download DOCX", disabled=True, use_container_width=True)
            else:
                try:
                    docx_data = generate_docx_bytes(st.session_state.template_bytes, text_data, image_data)
                    col_docx, col_pdf = st.columns(2)
                    with col_docx:
                        st.download_button(
                            label="Download DOCX", data=docx_data, file_name=get_download_filename(base_template_name, "docx"),
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            use_container_width=True, key="download_docx", on_click=purge_all_temporary_data
                        )
                    with col_pdf:
                        try:
                            pdf_data = docx_bytes_to_pdf(docx_data)
                        except Exception as e:
                            pdf_data = None
                            pdf_error = str(e)
                        if pdf_data is not None:
                            st.download_button(
                                label="Download PDF", data=pdf_data, file_name=get_download_filename(base_template_name, "pdf"),
                                mime="application/pdf", use_container_width=True, key="download_pdf",
                                on_click=purge_all_temporary_data
                            )
                        else:
                            st.button("Download PDF", disabled=True, use_container_width=True, key="download_pdf_disabled",
                                      help=f"PDF export unavailable: {pdf_error}")
                except Exception as e:
                    st.error(f"Error generating document: {str(e)}")
    else:
        st.info("Please upload or select a template to begin")
